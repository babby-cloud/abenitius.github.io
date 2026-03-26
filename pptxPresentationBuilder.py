import copy
import json
import math
import os
import sys
from datetime import datetime
from typing import Optional
from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.oxml.ns import qn
from pptx.util import Pt


# ── Slide index map ─────────────────────────────────────────────────────────────
# Maps logical variant names to their 0-based index in the template file.

SLIDE_INDICES = {
    # Start pages
    "start_page":                    1,
    "start_page_blue":               2,
    # Agenda
    "agenda":                        3,
    # Section headers
    "section_header":                5,
    "section_header_blue":           6,
    "section_header_basic":          7,
    # Title + text
    "title_text_small":              11,
    "title_text_medium":             12,
    "title_text_large":              13,
    "title_blue_line":               20,
    "title_center":                  21,
    "title_center_heading":          22,
    "title_center_heading_blue":     23,
    # Body only
    "body":                          24,
    "body_center":                   25,
    "body_center_blue":              26,
    # Case slides (text placeholders only; image placeholder left empty)
    "case_right_single":             31,
    "case_right_desktop":            32,
    "case_right_mobile":             33,
    "case_left_single":              37,
    "case_left_desktop":             38,
    "case_left_mobile":              39,
    # Quotes
    "quote_highlight":               46,
    "quote_highlight_blue":          47,
    "quote_body":                    49,
    "quote":                         50,
    # Columns
    "title_6_columns":               75,
    "title_2_columns":               76,
    "title_3_columns":               77,
    "title_4_columns":               78,
    # Text boxes
    "6_text_boxes_w_titles":         71,
    "6_text_boxes":                  72,
    "4_text_boxes_w_titles_img":     73,
    "3_text_boxes_w_titles":         74,
    # End
    "end_blue":                      146,
    "blank":                         147,
}

# Column placeholder indices per variant (list order = left to right)
COLUMN_PLACEHOLDER_INDICES = {
    "title_2_columns": [88, 90],
    "title_3_columns": [88, 89, 90],
    "title_4_columns": [88, 89, 90, 91],
    "title_6_columns": [88, 91, 114, 115, 116, 117],
}

# 6 text boxes w/ titles: (category_idx, heading_idx, body_idx) per box
# Visual order: Row 1 (top), Row 2 (middle), Row 3 (bottom)
TEXT_BOX_6_TITLES_INDICES = [
    (91,  90,  88),   # row 1, left
    (93,  92,  89),   # row 1, right
    (117, 116, 114),  # row 2, left
    (119, 118, 115),  # row 2, right
    (123, 122, 120),  # row 3, left
    (125, 124, 121),  # row 3, right
]

# 3 text boxes w/ titles: (title_idx, body_idx) per box
TEXT_BOX_3_TITLES_INDICES = [
    (125, 126),   # box left
    (129, 130),   # box center
    (127, 128),   # box right
]


# ── Core helpers ─────────────────────────────────────────────────────────────────

MAX_FONT_SIZE_PT = 96


def _estimate_fitting_font_size(text, width_emu, height_emu,
                                max_size_pt: float = 44, min_size_pt: float = 8,
                                margins=(91440, 91440, 45720, 45720)):
    """
    Find the largest whole-point font size (between *min_size_pt* and
    *max_size_pt*) that allows *text* to fit inside a text box whose
    dimensions are given in EMU.

    Uses heuristic character-width estimation (avg char ≈ 0.6 × font size)
    and word-wrap simulation.  Combined with PowerPoint's TEXT_TO_FIT_SHAPE
    auto-size this gives reliable results across slide types.

    margins: (left, right, top, bottom) in EMU — defaults match OOXML spec.
    """
    EMU_PER_PT = 12700
    SAFETY = 0.88
    usable_w = (width_emu - margins[0] - margins[1]) / EMU_PER_PT * SAFETY
    usable_h = (height_emu - margins[2] - margins[3]) / EMU_PER_PT * SAFETY

    if usable_w <= 0 or usable_h <= 0:
        return int(min_size_pt)

    segments = text.replace('\v', '\n').split('\n')

    for size in range(int(max_size_pt), int(min_size_pt) - 1, -1):
        chars_per_line = max(1, usable_w / (size * 0.6))
        total_lines = 0
        for segment in segments:
            words = segment.split()
            if not words:
                total_lines += 1
                continue
            line_len = 0
            seg_lines = 1
            for word in words:
                wl = len(word)
                if line_len == 0:
                    line_len = wl
                    if wl > chars_per_line:
                        extra = math.ceil(wl / chars_per_line) - 1
                        seg_lines += extra
                        line_len = wl - int(extra * chars_per_line)
                elif line_len + 1 + wl <= chars_per_line:
                    line_len += 1 + wl
                else:
                    seg_lines += 1
                    line_len = wl
                    if wl > chars_per_line:
                        extra = math.ceil(wl / chars_per_line) - 1
                        seg_lines += extra
                        line_len = wl - int(extra * chars_per_line)
            total_lines += seg_lines

        if total_lines * size * 1.2 <= usable_h:
            return size

    return int(min_size_pt)


class SlideBuilder:
    def __init__(self, template_path: str):
        if not os.path.exists(template_path):
            raise FileNotFoundError(f"Template not found: {template_path}")
        self.template_path = template_path
        self.template = Presentation(template_path)
        self.output = Presentation(template_path)
        self._remove_all_slides(self.output)
        self._ph_font_sizes: dict[int, float | None] = {}

    def _remove_all_slides(self, prs):
        # Remove slide references and their relationships
        r_ns = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
        sldIdLst = prs.slides._sldIdLst
        for sldId in list(sldIdLst):
            rId = sldId.get(f"{{{r_ns}}}id")
            prs.part.drop_rel(rId)
            sldIdLst.remove(sldId)

        # Remove slide sections (groups) stored in the presentation XML extension list
        prs_elem = prs.element
        ext_lst = prs_elem.find(qn("p:extLst"))
        if ext_lst is not None:
            for ext in list(ext_lst):
                uri = ext.get("uri", "")
                if uri == "{521415D9-36F7-43E2-AB2F-B90AF26B5E84}":
                    ext_lst.remove(ext)

    def _copy_slide(self, template_slide_index: int):
        if template_slide_index >= len(self.template.slides):
            raise IndexError(
                f"Template slide index {template_slide_index} out of range "
                f"(template has {len(self.template.slides)} slides)."
            )

        src_slide = self.template.slides[template_slide_index]

        # Match layout by name in the output presentation instead of always using [0]
        src_layout_name = src_slide.slide_layout.name
        tgt_layout = next(
            (l for l in self.output.slide_layouts if l.name == src_layout_name),
            self.output.slide_layouts[0]  # fallback only if no match found
        )

        new_slide = self.output.slides.add_slide(tgt_layout)

        # Replace shape tree
        dst_sp_tree = new_slide.shapes._spTree
        for child in list(dst_sp_tree):
            dst_sp_tree.remove(child)
        for child in src_slide.shapes._spTree:
            dst_sp_tree.append(copy.deepcopy(child))

        # Copy explicit slide background (if set directly on the slide)
        src_cSld = src_slide._element.find(qn("p:cSld"))
        dst_cSld = new_slide._element.find(qn("p:cSld"))
        if src_cSld is not None and dst_cSld is not None:
            src_bg = src_cSld.find(qn("p:bg"))
            if src_bg is not None:
                dst_bg = dst_cSld.find(qn("p:bg"))
                if dst_bg is not None:
                    dst_cSld.remove(dst_bg)
                dst_cSld.insert(0, copy.deepcopy(src_bg))

        # Capture font sizes from placeholder shapes before clearing
        self._ph_font_sizes = {}
        self._used_ph_indices: set[int] = set()
        for shape in new_slide.shapes:
            if shape.is_placeholder and shape.has_text_frame:
                ph_idx = shape.placeholder_format.idx
                size = self._resolve_max_font_size_pt(shape)
                if size is None:
                    size = self._resolve_font_size_from_layout(new_slide, ph_idx)
                self._ph_font_sizes[ph_idx] = size

        # Add empty "blocker" shapes for layout placeholders that have no
        # matching shape on the slide — prevents layout text (e.g. "CLICK TO
        # EDIT MASTER TEXT") from showing through.
        # Skip system placeholders (footer / date / slide-number / header)
        # because those are governed by presentation-level settings.
        _SYSTEM_PH_TYPES = frozenset(('ftr', 'dt', 'sldNum', 'hdr'))
        slide_ph_indices = {
            s.placeholder_format.idx
            for s in new_slide.shapes if s.is_placeholder
        }
        try:
            for ph in new_slide.slide_layout.placeholders:
                ph_idx = ph.placeholder_format.idx
                if ph_idx in slide_ph_indices:
                    continue
                ph_el = ph._element.find(f'.//{qn("p:ph")}')
                if ph_el is not None and ph_el.get('type', '') in _SYSTEM_PH_TYPES:
                    continue
                blocker = copy.deepcopy(ph._element)
                nvSpPr = blocker.find(qn('p:nvSpPr'))
                if nvSpPr is not None:
                    cNvPr = nvSpPr.find(qn('p:cNvPr'))
                    if cNvPr is not None:
                        cNvPr.set('id', str(900 + ph_idx))
                txBody = blocker.find(qn('p:txBody'))
                if txBody is not None:
                    for p in list(txBody.findall(qn('a:p'))):
                        txBody.remove(p)
                    empty_p = txBody.makeelement(qn('a:p'), {})
                    txBody.append(empty_p)
                new_slide.shapes._spTree.append(blocker)
        except Exception:
            pass

        # Clear placeholder text (filled shapes get overwritten by _set_text;
        # unused ones are removed later by _cleanup_slide)
        for shape in new_slide.shapes:
            if shape.is_placeholder and shape.has_text_frame:
                shape.text_frame.text = ""  # type: ignore[attr-defined]

        return new_slide

    def _resolve_max_font_size_pt(self, shape):
        """
        Walk the placeholder's text-formatting cascade to find the
        effective font size.  Returns size in points, or None when no
        explicit size can be found at any level.
        """
        if not shape.has_text_frame:
            return None

        tf = shape.text_frame

        for para in tf.paragraphs:
            for run in para.runs:
                if run.font.size is not None:
                    return run.font.size / 12700

        for para in tf.paragraphs:
            p = para._p
            pPr = p.find(qn('a:pPr'))
            if pPr is not None:
                defRPr = pPr.find(qn('a:defRPr'))
                if defRPr is not None and defRPr.get('sz'):
                    return int(defRPr.get('sz')) / 100
            endRPr = p.find(qn('a:endParaRPr'))
            if endRPr is not None and endRPr.get('sz'):
                return int(endRPr.get('sz')) / 100

        txBody = tf._txBody
        lstStyle = txBody.find(qn('a:lstStyle'))
        if lstStyle is not None:
            for lvl in range(1, 10):
                pPr = lstStyle.find(qn(f'a:lvl{lvl}pPr'))
                if pPr is not None:
                    defRPr = pPr.find(qn('a:defRPr'))
                    if defRPr is not None and defRPr.get('sz'):
                        return int(defRPr.get('sz')) / 100

        return None

    def _resolve_font_size_from_layout(self, slide, idx):
        """
        Walk the slide-layout → slide-master cascade looking for an
        explicit font size on the placeholder that matches *idx*.
        Returns size in points, or None.
        """
        for provider in (slide.slide_layout, slide.slide_layout.slide_master):
            try:
                for ph in provider.placeholders:
                    if ph.placeholder_format.idx == idx:
                        size = self._resolve_max_font_size_pt(ph)
                        if size is not None:
                            return size
            except Exception:
                continue
        return None

    def _calc_font_size(self, slide, idx: int, text: str) -> float:
        """
        Calculate the fitting font size (in pt) for *text* in placeholder
        *idx* without modifying the slide.  Returns the calculated size,
        or the capped template size as fallback.
        """
        for shape in slide.shapes:
            if shape.is_placeholder and shape.placeholder_format.idx == idx:
                if shape.has_text_frame:
                    tf = shape.text_frame
                    max_pt = min(self._ph_font_sizes.get(idx) or 44.0,
                                 MAX_FONT_SIZE_PT)

                    ml = tf.margin_left if tf.margin_left is not None else 91440
                    mr = tf.margin_right if tf.margin_right is not None else 91440
                    mt = tf.margin_top if tf.margin_top is not None else 45720
                    mb = tf.margin_bottom if tf.margin_bottom is not None else 45720

                    slide_w = self.output.slide_width or 12192000
                    slide_h = self.output.slide_height or 6858000
                    eff_w = min(shape.width, slide_w) if shape.width else 0
                    eff_h = min(shape.height, slide_h) if shape.height else 0

                    if eff_w and eff_h:
                        return _estimate_fitting_font_size(
                            text, eff_w, eff_h,
                            max_size_pt=max_pt, min_size_pt=8,
                            margins=(ml, mr, mt, mb),
                        )
                    return max_pt
        return 44.0

    def _set_text(self, slide, idx: int, text: str,
                  font_size_pt: float | None = None):
        """
        Set plain text on the placeholder identified by its idx,
        automatically scaling the font size to fit the text box.
        If *font_size_pt* is supplied it overrides the auto-calculation
        (used for uniform sizing across sibling placeholders).
        Silently skips image placeholders.
        """
        for shape in slide.shapes:
            if shape.is_placeholder and shape.placeholder_format.idx == idx:
                if shape.has_text_frame:
                    self._used_ph_indices.add(idx)
                    tf = shape.text_frame

                    tf.text = text

                    if font_size_pt is not None:
                        fit_pt: float | None = font_size_pt
                    else:
                        max_pt = min(self._ph_font_sizes.get(idx) or 44.0,
                                     MAX_FONT_SIZE_PT)
                        ml = tf.margin_left if tf.margin_left is not None else 91440
                        mr = tf.margin_right if tf.margin_right is not None else 91440
                        mt = tf.margin_top if tf.margin_top is not None else 45720
                        mb = tf.margin_bottom if tf.margin_bottom is not None else 45720

                        slide_w = self.output.slide_width or 12192000
                        slide_h = self.output.slide_height or 6858000
                        eff_w = min(shape.width, slide_w) if shape.width else 0
                        eff_h = min(shape.height, slide_h) if shape.height else 0

                        if eff_w and eff_h:
                            fit_pt = _estimate_fitting_font_size(
                                text, eff_w, eff_h,
                                max_size_pt=max_pt, min_size_pt=8,
                                margins=(ml, mr, mt, mb),
                            )
                        else:
                            fit_pt = None

                    if fit_pt is not None:
                        size_emu = Pt(fit_pt)
                        for para in tf.paragraphs:
                            for run in para.runs:
                                run.font.size = size_emu

                    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                return

    def _cleanup_slide(self, slide):
        """
        Remove placeholder shapes that were never filled by _set_text.
        This eliminates empty dashed boxes and PowerPoint's ghost prompt
        text (e.g. Lorem ipsum from the layout) on unused placeholders.
        """
        removals = []
        for shape in slide.shapes:
            if not shape.is_placeholder:
                continue
            if shape.placeholder_format.idx in self._used_ph_indices:
                continue
            removals.append(shape._element)
        sp_tree = slide.shapes._spTree
        for elem in removals:
            sp_tree.remove(elem)

    def save(self, output_path: str):
        props = self.output.core_properties
        now = datetime.now()
        props.created = now
        props.modified = now
        props.author = "iO Presentation Generator"
        props.last_modified_by = "iO Presentation Generator"
        props.revision = 1

        self.output.save(output_path)
        print(f"Saved to {output_path}  ({len(self.output.slides)} slides)")


# ── Slide functions ───────────────────────────────────────────────────────────────

def add_start_page(
    builder: SlideBuilder,
    title: str,
    subject: str,
    presenter: str,
    blue_line: bool = False,
):
    """
    Opening / title slide.
    idx 0   → title (large, bottom area)
    idx 114 → subject (top-left label)
    idx 112 → presenter name / title / company (top-right)
    """
    variant = "start_page_blue" if blue_line else "start_page"
    slide = builder._copy_slide(SLIDE_INDICES[variant])
    builder._set_text(slide, 0,   title)
    builder._set_text(slide, 114, subject)
    builder._set_text(slide, 112, presenter)
    return slide


def add_section_header(
    builder: SlideBuilder,
    title: str,
    subheading: Optional[str] = None,
    variant: str = "default",
):
    """
    Section divider slide.
    variant: "default" | "blue_line" | "basic"
    idx 0   → section title
    idx 114 → subheading (not available in "basic")
    """
    key_map = {
        "default":   "section_header",
        "blue_line": "section_header_blue",
        "basic":     "section_header_basic",
    }
    if variant not in key_map:
        raise ValueError(f"variant must be one of {list(key_map.keys())}")

    slide = builder._copy_slide(SLIDE_INDICES[key_map[variant]])
    builder._set_text(slide, 0, title)
    if subheading and variant != "basic":
        builder._set_text(slide, 114, subheading)
    return slide


def add_agenda(
    builder: SlideBuilder,
    sections: list[str],
):
    """
    Agenda slide with a list of section titles.
    The sections are joined with vertical tab (\\v) which renders as
    line breaks within the same text box — matching the template's style.
    idx 0 → full agenda text block
    """
    slide = builder._copy_slide(SLIDE_INDICES["agenda"])
    builder._set_text(slide, 0, "\v".join(sections))
    return slide


def add_title_text(
    builder: SlideBuilder,
    title: str,
    body: str,
    subject: Optional[str] = None,
    size: str = "medium",
    centered: bool = False,
    subheading: Optional[str] = None,
    blue_line: bool = False,
):
    """
    Workhorse content slide: title + body text.
    size:     "small" | "medium" | "large"   (ignored when centered=True)
    centered: use the centered title layout
    subheading: only used when centered=True
    blue_line:  use blue line accent variant

    Placeholder mapping:
      idx 0   → title
      idx 113 → subject label (top-left)
      idx 114 → body (size variants) OR subheading (centered variants)
      idx 12  → body (title_blue_line layout)
    """
    if centered:
        if blue_line:
            key = "title_center_heading_blue"
        elif subheading:
            key = "title_center_heading"
        else:
            key = "title_center"
    elif blue_line:
        key = "title_blue_line"
    else:
        size_map = {"small": "title_text_small", "medium": "title_text_medium", "large": "title_text_large"}
        if size not in size_map:
            raise ValueError(f"size must be one of {list(size_map.keys())}")
        key = size_map[size]

    slide = builder._copy_slide(SLIDE_INDICES[key])
    builder._set_text(slide, 0, title)

    if subject:
        builder._set_text(slide, 113, subject)

    if key == "title_blue_line":
        builder._set_text(slide, 12, body)
    elif key in ("title_center_heading", "title_center_heading_blue"):
        builder._set_text(slide, 114, subheading or "")
    else:
        builder._set_text(slide, 114, body)

    return slide


def add_body_only(
    builder: SlideBuilder,
    body: str,
    subject: Optional[str] = None,
    centered: bool = False,
    blue_line: bool = False,
):
    """
    Body-only slide (no title).
    idx 113 → subject label
    idx 12  → body text
    """
    if centered and blue_line:
        key = "body_center_blue"
    elif centered:
        key = "body_center"
    else:
        key = "body"

    slide = builder._copy_slide(SLIDE_INDICES[key])
    if subject:
        builder._set_text(slide, 113, subject)
    builder._set_text(slide, 12, body)
    return slide


def add_columns(
    builder: SlideBuilder,
    title: str,
    columns: list[dict],
    subject: Optional[str] = None,
):
    """
    Multi-column slide. Number of items in `columns` determines layout (2–6).
    Each column: { "heading": str, "body": str }
    Combined as "heading\\n\\nbody" into each column placeholder.

    idx 0   → title
    idx 113 → subject label
    Column placeholders vary by count — see COLUMN_PLACEHOLDER_INDICES.
    """
    count = len(columns)
    key_map = {2: "title_2_columns", 3: "title_3_columns", 4: "title_4_columns", 6: "title_6_columns"}
    if count not in key_map:
        raise ValueError(f"columns must have 2, 3, 4, or 6 items — got {count}")

    key = key_map[count]
    slide = builder._copy_slide(SLIDE_INDICES[key])
    builder._set_text(slide, 0, title)
    if subject:
        builder._set_text(slide, 113, subject)

    col_texts = [
        f"{columns[i].get('heading', '')}\n\n{columns[i].get('body', '')}"
        for i in range(count)
    ]
    sizes = [
        builder._calc_font_size(slide, col_idx, text)
        for col_idx, text in zip(COLUMN_PLACEHOLDER_INDICES[key], col_texts)
        if text.strip()
    ]
    uniform = min(sizes) if sizes else None

    for col_idx, text in zip(COLUMN_PLACEHOLDER_INDICES[key], col_texts):
        builder._set_text(slide, col_idx, text, font_size_pt=uniform)

    return slide


def add_text_boxes(
    builder: SlideBuilder,
    title: str,
    boxes: list[dict],
    subject: Optional[str] = None,
):
    """
    Text box grid slide. Supports 3 or 6 boxes.
    Each box: { "category": str, "heading": str, "body": str }
    (category and heading are optional for 3-box layout)

    3 boxes → layout "3 text boxes w/ titles"   (slide 74)
    6 boxes → layout "6 text boxes w/ titles"   (slide 71)

    idx 0   → title
    idx 113 → subject label
    idx 12  → intro/body text above boxes (3-box layout only); pass as title if needed
    """
    count = len(boxes)
    if count == 3:
        slide = builder._copy_slide(SLIDE_INDICES["3_text_boxes_w_titles"])
        builder._set_text(slide, 0, title)
        if subject:
            builder._set_text(slide, 113, subject)

        heading_sizes: list[float] = []
        body_sizes: list[float] = []
        for i, (title_idx, body_idx) in enumerate(TEXT_BOX_3_TITLES_INDICES):
            if i < len(boxes):
                h = boxes[i].get("heading", "")
                b = boxes[i].get("body", "")
                if h.strip():
                    heading_sizes.append(builder._calc_font_size(slide, title_idx, h))
                if b.strip():
                    body_sizes.append(builder._calc_font_size(slide, body_idx, b))
        uniform_h = min(heading_sizes) if heading_sizes else None
        uniform_b = min(body_sizes) if body_sizes else None

        for i, (title_idx, body_idx) in enumerate(TEXT_BOX_3_TITLES_INDICES):
            if i < len(boxes):
                builder._set_text(slide, title_idx, boxes[i].get("heading", ""),
                                  font_size_pt=uniform_h)
                builder._set_text(slide, body_idx, boxes[i].get("body", ""),
                                  font_size_pt=uniform_b)

    elif count == 6:
        slide = builder._copy_slide(SLIDE_INDICES["6_text_boxes_w_titles"])
        builder._set_text(slide, 0, title)
        if subject:
            builder._set_text(slide, 113, subject)
        indices = TEXT_BOX_6_TITLES_INDICES

        cat_sizes: list[float] = []
        hdg_sizes: list[float] = []
        bdy_sizes: list[float] = []
        for i, (cat_idx, hdg_idx, bdy_idx) in enumerate(indices):
            if i < len(boxes):
                c = boxes[i].get("category", "")
                h = boxes[i].get("heading", "")
                b = boxes[i].get("body", "")
                if cat_idx and c.strip():
                    cat_sizes.append(builder._calc_font_size(slide, cat_idx, c))
                if hdg_idx and h.strip():
                    hdg_sizes.append(builder._calc_font_size(slide, hdg_idx, h))
                if bdy_idx and b.strip():
                    bdy_sizes.append(builder._calc_font_size(slide, bdy_idx, b))

        uniform_c = min(cat_sizes) if cat_sizes else None
        uniform_hd = min(hdg_sizes) if hdg_sizes else None
        uniform_bd = min(bdy_sizes) if bdy_sizes else None

        for i, (cat_idx, hdg_idx, bdy_idx) in enumerate(indices):
            if i < len(boxes):
                if cat_idx:
                    builder._set_text(slide, cat_idx, boxes[i].get("category", ""),
                                      font_size_pt=uniform_c)
                if hdg_idx:
                    builder._set_text(slide, hdg_idx, boxes[i].get("heading", ""),
                                      font_size_pt=uniform_hd)
                if bdy_idx:
                    builder._set_text(slide, bdy_idx, boxes[i].get("body", ""),
                                      font_size_pt=uniform_bd)
    else:
        raise ValueError(f"boxes must have 3 or 6 items — got {count}")

    return slide


def add_quote(
    builder: SlideBuilder,
    quote: str,
    attribution: Optional[str] = None,
    subject: Optional[str] = None,
    highlight: bool = False,
    body: Optional[str] = None,
    blue_line: bool = False,
):
    """
    Quote slide.
    highlight:  use the highlight variant (large quote, no attribution block)
    body:       adds a body text block alongside the quote ("quote + body" layout)
    blue_line:  blue accent line (highlight variant only)

    idx 88  → quote text
    idx 112 → attribution (Name / Title / Company)  [quote layout only]
    idx 113 → subject label
    idx 0   → title  [quote_body layout only]
    """
    if body:
        key = "quote_body"
    elif highlight and blue_line:
        key = "quote_highlight_blue"
    elif highlight:
        key = "quote_highlight"
    else:
        key = "quote"

    slide = builder._copy_slide(SLIDE_INDICES[key])

    if subject:
        builder._set_text(slide, 113, subject)

    if key == "quote_body":
        builder._set_text(slide, 88, quote)
        builder._set_text(slide, 0,  body or "")
    else:
        builder._set_text(slide, 88, quote)
        if attribution and key == "quote":
            builder._set_text(slide, 112, attribution)

    return slide


def add_case_slide(
    builder: SlideBuilder,
    title: str,
    body: str,
    client: Optional[str] = None,
    subject: Optional[str] = None,
    image_side: str = "right",
    device: str = "desktop",
):
    """
    Case study slide. Image placeholder is left empty for manual insertion.
    image_side: "left" | "right"
    device:     "desktop" | "mobile" | "single"

    idx 0   → case headline
    idx 84  → description body
    idx 113 → subject label
    idx 114 → client name (bottom)
    """
    key_map = {
        ("right", "single"):  "case_right_single",
        ("right", "desktop"): "case_right_desktop",
        ("right", "mobile"):  "case_right_mobile",
        ("left",  "single"):  "case_left_single",
        ("left",  "desktop"): "case_left_desktop",
        ("left",  "mobile"):  "case_left_mobile",
    }
    key = key_map.get((image_side, device))
    if not key:
        raise ValueError(f"Invalid combination: image_side='{image_side}', device='{device}'")

    slide = builder._copy_slide(SLIDE_INDICES[key])
    builder._set_text(slide, 0,   title)
    builder._set_text(slide, 84,  body)
    if subject:
        builder._set_text(slide, 113, subject)
    if client:
        builder._set_text(slide, 114, client)
    return slide


def add_end_slide(
    builder: SlideBuilder,
    tagline: Optional[str] = None,
    presenter: Optional[str] = None,
    blue_line: bool = True,
):
    """
    Closing slide.
    idx 0   → tagline (large centered text)
    idx 115 → presenter name (top-right corner)
    idx 112 → logo/empty area (left untouched)
    """
    key = "end_blue" if blue_line else "blank"
    slide = builder._copy_slide(SLIDE_INDICES[key])
    if tagline:
        builder._set_text(slide, 0, tagline)
    if presenter:
        builder._set_text(slide, 115, presenter)
    return slide


# ── Build from JSON ───────────────────────────────────────────────────────────────

FUNCTION_MAP = {
    "add_start_page":    add_start_page,
    "add_section_header": add_section_header,
    "add_agenda":        add_agenda,
    "add_title_text":    add_title_text,
    "add_body_only":     add_body_only,
    "add_columns":       add_columns,
    "add_text_boxes":    add_text_boxes,
    "add_quote":         add_quote,
    "add_case_slide":    add_case_slide,
    "add_end_slide":     add_end_slide,
}


def build_from_instructions(
    instructions: list[dict], template_path: str, output_path: str,
):
    """Build a presentation from an in-memory instruction list."""
    builder = SlideBuilder(template_path)

    for i, step in enumerate(instructions):
        fn_name = step.get("function")
        params  = step.get("params", {})

        if fn_name not in FUNCTION_MAP:
            raise ValueError(
                f"Step {i}: unknown function '{fn_name}'. "
                f"Allowed: {sorted(FUNCTION_MAP.keys())}"
            )

        try:
            slide = FUNCTION_MAP[fn_name](builder, **params)
            if slide is not None:
                builder._cleanup_slide(slide)
        except TypeError as e:
            raise TypeError(f"Step {i} ({fn_name}): invalid params — {e}") from e

    builder.save(output_path)


def build_from_json(instructions_path: str, template_path: str, output_path: str):
    """Load a JSON instructions file and build the presentation."""
    with open(instructions_path, "r", encoding="utf-8") as f:
        instructions = json.load(f)
    build_from_instructions(instructions, template_path, output_path)


# ── Entry point ───────────────────────────────────────────────────────────────────

if __name__ == "__main__":
    if len(sys.argv) < 4:
        print("Usage: python slide_builder.py <instructions.json> <template.pptx> <output.pptx>")
        sys.exit(1)

    build_from_json(
        instructions_path=sys.argv[1],
        template_path=sys.argv[2],
        output_path=sys.argv[3],
    )